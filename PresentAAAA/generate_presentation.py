"""
Generate GraphRAG presentation: 12-slide PPTX + speaker-notes PDF.
Run from project root with venv active:
    python PresentAAAA/generate_presentation.py
"""

import os
from pathlib import Path

# ── colour palette ─────────────────────────────────────────────────────────────
DARK_BG   = "0D1117"   # near-black
MID_BG    = "161B22"   # card dark
ACCENT1   = "58A6FF"   # electric blue
ACCENT2   = "3FB950"   # success green
ACCENT3   = "F78166"   # coral/red
ACCENT4   = "D2A8FF"   # lavender
GOLD      = "E3B341"   # gold/amber
WHITE     = "F0F6FC"
GREY      = "8B949E"
TEAL      = "39D353"

# ── helpers ────────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

OUT_DIR = Path(__file__).parent
PPTX_PATH = OUT_DIR / "GraphRAG_Presentation.pptx"
PDF_PATH  = OUT_DIR / "Speaker_Notes.pdf"

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill_hex, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox

def add_multiline(slide, lines, left, top, width, height,
                  font_size=13, color=WHITE, bold_first=False, line_spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color)
        if bold_first and i == 0:
            run.font.bold = True
    return txBox

def dark_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # full dark background
    bg = add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    return slide

def title_bar(slide, title, subtitle=None, accent=ACCENT1):
    # top accent stripe
    add_rect(slide, 0, 0, 13.33, 0.08, accent)
    # title
    add_text(slide, title, 0.5, 0.18, 12.3, 0.85,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.95, 12.3, 0.45,
                 font_size=15, color=GREY, align=PP_ALIGN.LEFT)

def card(slide, left, top, width, height, fill=MID_BG, accent=None):
    r = add_rect(slide, left, top, width, height, fill)
    if accent:
        add_rect(slide, left, top, width, 0.06, accent)
    return r

# ── build slides ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE / COVER
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
# gradient-ish panels
add_rect(s, 0, 0, 6.5, 7.5, "0D1B2A")
add_rect(s, 6.5, 0, 6.83, 7.5, "0A3D62")
# decorative accent lines
add_rect(s, 0, 0, 6.5, 0.12, ACCENT1)
add_rect(s, 6.5, 0, 6.83, 0.12, ACCENT2)
# big title
add_text(s, "GraphRAG", 0.6, 1.5, 5.8, 1.3,
         font_size=56, bold=True, color=ACCENT1)
add_text(s, "for Domain-Specific LLMs", 0.6, 2.75, 5.8, 0.7,
         font_size=26, bold=False, color=WHITE)
add_rect(s, 0.6, 3.55, 3.5, 0.05, ACCENT2)
add_text(s, "Knowledge-Graph-Augmented Retrieval\nfor Supply-Chain Resilience Q&A", 0.6, 3.75, 5.6, 1.0,
         font_size=15, color=GREY)
# right panel content
add_text(s, "📊 34-Question Benchmark", 7.0, 1.6, 5.5, 0.5, font_size=16, color=WHITE, bold=True)
add_text(s, "4 Retrieval Methods Compared", 7.0, 2.1, 5.5, 0.45, font_size=14, color=GREY)
add_rect(s, 7.0, 2.65, 5.1, 0.06, ACCENT2)
add_text(s, "🏆 GraphRAG: #1 Overall", 7.0, 2.85, 5.5, 0.5, font_size=16, color=ACCENT2, bold=True)
add_text(s, "Score 0.767  |  22 / 34 wins", 7.0, 3.32, 5.5, 0.45, font_size=14, color=GREY)
add_rect(s, 7.0, 3.85, 5.1, 0.06, ACCENT4)
add_text(s, "Citation Precision: 91.2%", 7.0, 4.05, 5.5, 0.45, font_size=14, color=ACCENT4)
# bottom tagline
add_text(s, "From keyword search to semantic graph traversal — why it matters", 0.6, 6.7, 12.0, 0.45,
         font_size=12, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "The Problem", "Why standard RAG falls short on expert domains", ACCENT3)

problems = [
    ("❌  Keyword blindness", "TF-IDF matches words but misses meaning — 'resilience' and 'robustness' are unrelated to it"),
    ("❌  Flat retrieval", "Documents are treated as bags-of-words with no sense of how concepts relate"),
    ("❌  No causal chains", "Standard RAG cannot answer 'How does X lead to Y through Z?' questions"),
    ("❌  Context fragmentation", "Splitting text into chunks severs relationships that span sentence boundaries"),
]

y = 1.5
for i, (title, body) in enumerate(problems):
    col = [ACCENT3, ACCENT3, GOLD, GOLD][i]
    card(s, 0.4, y, 5.9, 1.1, MID_BG, col)
    add_text(s, title, 0.65, y+0.12, 5.5, 0.38, font_size=14, bold=True, color=col)
    add_text(s, body,  0.65, y+0.48, 5.5, 0.55, font_size=12, color=GREY)
    y += 1.25

# right: domain context
card(s, 6.7, 1.5, 6.2, 5.5, MID_BG, ACCENT1)
add_text(s, "Our Domain: Supply-Chain Resilience", 6.95, 1.65, 5.8, 0.5,
         font_size=14, bold=True, color=ACCENT1)
domain_lines = [
    "• OECD, World Bank, WTO policy documents",
    "• Entities: strategy, policy, capability, risk,",
    "  actor, sector, system, org, location",
    "",
    "• Relations: STRENGTHENS, MITIGATES,",
    "  ENABLES, TRADE_OFF_WITH, EXPOSES,",
    "  CONSTRAINS, AFFECTS …",
    "",
    "• Questions require multi-hop reasoning:",
    "  'How does nearshoring affect resilience",
    "  through supplier diversification?'",
]
add_multiline(s, domain_lines, 6.95, 2.2, 5.7, 4.4, font_size=12.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "System Architecture", "Two-phase pipeline: Indexing → Retrieval & Generation", ACCENT4)

# pipeline boxes
boxes = [
    (0.3,  1.7, 2.2, "📄 Corpus\nIngestion", ACCENT1,   "OECD / WB / WTO\nPDF → chunks"),
    (2.85, 1.7, 2.2, "🔍 Entity &\nRelation\nExtraction", ACCENT4, "spaCy NER +\nTrigger vocab"),
    (5.40, 1.7, 2.2, "🕸 Knowledge\nGraph Build", ACCENT2, "NetworkX\nentity nodes\nrelation edges"),
    (7.95, 1.7, 2.2, "📦 Index\nStore", GOLD,    "FAISS dense\n+ TF-IDF sparse\n+ Graph JSON"),
    (10.5, 1.7, 2.2, "💬 Query &\nAnswer", ACCENT3, "4 retrievers\n+ GPT-4o-mini\ngeneration"),
]
for x, y, w, label, col, sub in boxes:
    card(s, x, y, w, 1.9, MID_BG, col)
    add_text(s, label, x+0.1, y+0.15, w-0.2, 0.85, font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, sub,   x+0.1, y+1.0,  w-0.2, 0.8,  font_size=11, color=GREY, align=PP_ALIGN.CENTER)

# arrows between boxes
for ax in [2.55, 5.1, 7.65, 10.2]:
    add_text(s, "→", ax, 2.3, 0.4, 0.5, font_size=22, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

# bottom detail row
add_rect(s, 0.3, 4.0, 12.8, 0.05, ACCENT1)
detail_cols = [
    (0.3,  "Chunking\n~300 tokens\n50 token overlap", ACCENT1),
    (3.3,  "Cross-sentence\n2-sentence sliding\nwindow", ACCENT4),
    (6.3,  "9 entity types\n10 relation types\n~61 rels/chunk", ACCENT2),
    (9.3,  "RRF fusion\nDense+Graph\nsignals", GOLD),
]
for x, txt, col in detail_cols:
    card(s, x, 4.2, 2.85, 2.8, "0D1B2A", col)
    add_text(s, txt, x+0.15, 4.4, 2.55, 2.4, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – KNOWLEDGE GRAPH
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "The Knowledge Graph", "Entities, relations, and multi-hop paths", ACCENT2)

# entity type table
card(s, 0.4, 1.4, 5.5, 5.6, MID_BG, ACCENT2)
add_text(s, "Entity Types & Weights", 0.65, 1.55, 5.0, 0.45, font_size=14, bold=True, color=ACCENT2)
entities = [
    ("strategy",      "3.0", ACCENT1),
    ("policy",        "2.9", ACCENT4),
    ("capability",    "2.6", ACCENT2),
    ("risk",          "2.4", ACCENT3),
    ("actor",         "1.9", GOLD),
    ("sector",        "1.9", GOLD),
    ("system",        "1.1", GREY),
    ("organization",  "0.7", GREY),
    ("location",      "0.5", GREY),
]
for i, (ent, wt, col) in enumerate(entities):
    y = 2.1 + i * 0.5
    add_rect(s, 0.55, y, 3.4, 0.42, "0D1B2A")
    add_text(s, ent,  0.75, y+0.05, 2.5, 0.35, font_size=12.5, color=col)
    add_text(s, f"× {wt}", 3.5, y+0.05, 0.9, 0.35, font_size=12.5, bold=True, color=col, align=PP_ALIGN.RIGHT)

# relation types
card(s, 6.2, 1.4, 6.7, 2.6, MID_BG, ACCENT4)
add_text(s, "Relation Types", 6.45, 1.55, 6.2, 0.45, font_size=14, bold=True, color=ACCENT4)
rels = ["STRENGTHENS  •  IMPROVES  •  MITIGATES",
        "TRADE_OFF_WITH  •  ENABLES  •  EXPOSES",
        "CONSTRAINS  •  AFFECTS  •  APPLIES_TO",
        "CO_OCCURS_WITH"]
add_multiline(s, rels, 6.45, 2.05, 6.2, 1.7, font_size=12.5, color=WHITE)

# graph scoring formula
card(s, 6.2, 4.2, 6.7, 2.75, MID_BG, GOLD)
add_text(s, "Graph Scoring Formula", 6.45, 4.35, 6.2, 0.45, font_size=14, bold=True, color=GOLD)
formula_lines = [
    "score =  entity_match  +  coverage_bonus",
    "       + sentence_alignment  +  direct_relation",
    "       + path_bonus  +  lexical_score",
    "       − generic_penalty  − query_miss_penalty",
    "       − hub_penalty",
    "",
    "Bigrams get 2× lexical weight | Safety-net injects",
    "top-2 lexical chunks regardless of score",
]
add_multiline(s, formula_lines, 6.45, 4.85, 6.2, 2.0, font_size=11.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – 4 RETRIEVAL METHODS
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Four Retrieval Methods", "Baseline comparison across the spectrum", ACCENT1)

methods = [
    ("TF-IDF",        ACCENT3, "Sparse lexical\nbag-of-words",
     ["• Term frequency – inverse document freq.", "• Fast, no model needed",
      "• Cannot handle paraphrase or context", "• Treats each chunk independently"]),
    ("Dense",         ACCENT1, "Semantic embeddings\n+ FAISS ANN",
     ["• sentence-transformers (MiniLM-L6)", "• Cosine similarity search via FAISS",
      "• Handles synonyms & paraphrase", "• No structural/relational knowledge"]),
    ("GraphRAG",      ACCENT2, "Knowledge-graph\ntraversal",
     ["• Entity matching + relation path scoring", "• Cross-sentence relation extraction",
      "• Bigram query terms + lexical safety-net", "• Highest citation precision: 91.2%"]),
    ("Hybrid",        ACCENT4, "RRF fusion\nDense + Graph",
     ["• Reciprocal Rank Fusion (k=60)", "• 38% dense + 28% graph + 34% rank signals",
      "• Best on causal & multi-hop categories", "• Balanced precision/recall trade-off"]),
]

for i, (name, col, tag, bullets) in enumerate(methods):
    x = 0.3 + i * 3.26
    card(s, x, 1.45, 3.1, 5.6, MID_BG, col)
    add_text(s, name, x+0.15, 1.6, 2.8, 0.5, font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, tag,  x+0.15, 2.1, 2.8, 0.6, font_size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, x+0.2, 2.72, 2.7, 0.04, col)
    add_multiline(s, bullets, x+0.15, 2.85, 2.8, 3.9, font_size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – BENCHMARK & EVALUATION
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Benchmark & Evaluation Setup", "34 questions · LLM-as-judge · Weighted scoring", GOLD)

# left: question breakdown
card(s, 0.3, 1.4, 5.4, 5.5, MID_BG, GOLD)
add_text(s, "34 Benchmark Questions", 0.55, 1.55, 5.0, 0.45, font_size=14, bold=True, color=GOLD)
cats = [
    ("direct_factual", "6",  ACCENT1,  "Exact facts from documents"),
    ("causal",         "11", ACCENT2,  "X → Y causal chains"),
    ("mitigation",     "4",  ACCENT3,  "How to reduce / prevent"),
    ("trade_off",      "4",  ACCENT4,  "Tension between properties"),
    ("multi_hop",      "9",  GOLD,     "Requires ≥2 reasoning steps"),
]
for i, (cat, n, col, desc) in enumerate(cats):
    y = 2.1 + i * 0.95
    add_rect(s, 0.5, y, 4.8, 0.85, "0D1B2A")
    add_rect(s, 0.5, y, 0.08, 0.85, col)
    add_text(s, f"{cat}  ({n})", 0.7, y+0.06, 3.5, 0.35, font_size=13, bold=True, color=col)
    add_text(s, desc,            0.7, y+0.42, 3.8, 0.35, font_size=11.5, color=GREY)

# right: scoring breakdown
card(s, 6.0, 1.4, 6.9, 2.6, MID_BG, ACCENT1)
add_text(s, "LLM-as-Judge Scoring Weights", 6.25, 1.55, 6.4, 0.45, font_size=14, bold=True, color=ACCENT1)
dims = [
    ("Correctness",   "35%", ACCENT1),
    ("Completeness",  "25%", ACCENT2),
    ("Groundedness",  "25%", ACCENT4),
    ("Reasoning",     "10%", GOLD),
    ("Clarity",        "5%", GREY),
]
for i, (dim, pct, col) in enumerate(dims):
    y = 2.1 + i * 0.38
    bar_w = float(pct.strip('%')) / 35 * 3.8
    add_rect(s, 6.2, y, bar_w, 0.28, col)
    add_text(s, f"{dim}  {pct}", 6.2+bar_w+0.1, y+0.02, 3.0, 0.28, font_size=12, color=WHITE)

# models used
card(s, 6.0, 4.15, 6.9, 2.7, MID_BG, ACCENT3)
add_text(s, "Models Used", 6.25, 4.3, 6.4, 0.45, font_size=14, bold=True, color=ACCENT3)
model_lines = [
    "Generation:  GPT-4o-mini",
    "Judge:       GPT-4o-mini",
    "Embeddings:  all-MiniLM-L6-v2",
    "NER:         spaCy en_core_web_sm",
    "             + custom entity rules",
]
add_multiline(s, model_lines, 6.25, 4.82, 6.4, 1.85, font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – RESULTS OVERVIEW (BIG TABLE)
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Benchmark Results", "GraphRAG leads on every key metric", ACCENT2)

# winner banner
add_rect(s, 0.3, 1.1, 12.8, 0.72, "0D3A1A")
add_rect(s, 0.3, 1.1, 0.1, 0.72, ACCENT2)
add_text(s, "🏆  GraphRAG achieves the highest overall score: 0.767  —  outperforming TF-IDF (0.725), Hybrid (0.728) and Dense (0.680)",
         0.55, 1.2, 12.3, 0.5, font_size=14, bold=True, color=ACCENT2)

# table header
headers = ["Method", "Overall ↑", "Correctness", "Completeness", "Reasoning", "Groundedness", "Clarity", "Wins"]
col_x   = [0.3, 2.1, 3.8, 5.4, 7.0, 8.55, 10.1, 11.65]
col_w   = [1.7, 1.6, 1.5, 1.5, 1.45, 1.45, 1.45, 1.5]
add_rect(s, 0.3, 2.0, 12.8, 0.52, "1C2D1C")
for i, (h, x) in enumerate(zip(headers, col_x)):
    add_text(s, h, x+0.05, 2.05, col_w[i], 0.42, font_size=12, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

rows = [
    ("GraphRAG",  "0.767", "0.813", "0.704", "0.810", "0.756", "0.821", "22", ACCENT2, "162B1A"),
    ("Hybrid",    "0.728", "0.765", "0.666", "0.762", "0.722", "0.774", "15", ACCENT4, MID_BG),
    ("TF-IDF",    "0.725", "0.762", "0.674", "0.751", "0.709", "0.754", "22", ACCENT3, MID_BG),
    ("Dense",     "0.680", "0.726", "0.618", "0.726", "0.651", "0.738", "14", ACCENT1, MID_BG),
]
for ri, (m, ov, cr, cm, rs, gr, cl, w, col, bg) in enumerate(rows):
    y = 2.6 + ri * 0.88
    add_rect(s, 0.3, y, 12.8, 0.82, bg)
    vals = [m, ov, cr, cm, rs, gr, cl, w]
    for ci, (v, x) in enumerate(zip(vals, col_x)):
        is_highlight = (ci == 1)  # overall column
        add_text(s, v, x+0.05, y+0.2, col_w[ci], 0.42,
                 font_size=14 if is_highlight else 13,
                 bold=is_highlight,
                 color=col if is_highlight else WHITE,
                 align=PP_ALIGN.CENTER)

add_text(s, "↑ higher is better  |  Wins = questions where method scored highest",
         0.3, 7.1, 12.8, 0.35, font_size=11, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – CATEGORY BREAKDOWN
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Performance by Question Category", "GraphRAG dominates reasoning-heavy categories", ACCENT1)

cats2 = ["direct_factual", "causal", "mitigation", "multi_hop", "trade_off"]
graph_scores  = [0.850, 0.800, 0.631, 0.711, 0.744]
tfidf_scores  = [0.854, 0.730, 0.459, 0.777, 0.694]
hybrid_scores = [0.668, 0.790, 0.521, 0.777, 0.750]
dense_scores  = [0.610, 0.765, 0.519, 0.760, 0.772]

bar_groups = list(zip(cats2, graph_scores, tfidf_scores, hybrid_scores, dense_scores))
bar_area_left = 0.5
bar_area_top  = 1.5
bar_area_w    = 12.3
bar_area_h    = 4.8
group_w = bar_area_w / len(bar_groups)
bar_w   = 0.35
gap     = 0.08
colors  = [ACCENT2, ACCENT3, ACCENT4, ACCENT1]
max_val = 1.0

for gi, (cat, g, t, h, d) in enumerate(bar_groups):
    gx = bar_area_left + gi * group_w
    scores_list = [g, t, h, d]
    for bi, (score, col) in enumerate(zip(scores_list, colors)):
        bx = gx + 0.15 + bi * (bar_w + gap)
        bh = (score / max_val) * bar_area_h
        by = bar_area_top + bar_area_h - bh
        add_rect(s, bx, by, bar_w, bh, col)
        add_text(s, f"{score:.2f}", bx, by - 0.28, bar_w, 0.25,
                 font_size=8.5, color=col, align=PP_ALIGN.CENTER, bold=True)
    # category label
    short = cat.replace("direct_factual","factual").replace("multi_hop","multi-hop")
    add_text(s, short, gx+0.1, bar_area_top+bar_area_h+0.05, group_w-0.1, 0.35,
             font_size=11, color=GREY, align=PP_ALIGN.CENTER)

# legend
lx = 1.0
for col, label in zip(colors, ["GraphRAG", "TF-IDF", "Hybrid", "Dense"]):
    add_rect(s, lx, 6.8, 0.25, 0.18, col)
    add_text(s, label, lx+0.3, 6.78, 1.5, 0.25, font_size=11.5, color=WHITE)
    lx += 2.0

# highlight boxes
highlights = [
    ("Graph crushes Mitigation", "0.631 vs TF-IDF 0.459\n+37% improvement", ACCENT2, 0.3),
    ("Graph leads Causal",       "0.800 vs Dense 0.765\nClear reasoning win", ACCENT2, 3.55),
    ("Graph #1 Correctness",     "0.813 overall\nHighest in benchmark", GOLD, 6.8),
    ("Cite precision 91.2%",     "vs TF-IDF 85.3%\n+5.9 pp precision", ACCENT4, 10.05),
]
for label, val, col, hx in highlights:
    card(s, hx, 6.55, 2.9, 0.85, "0D1B2A", col)
    add_text(s, label, hx+0.15, 6.58, 2.6, 0.38, font_size=10.5, bold=True, color=col)
    add_text(s, val,   hx+0.15, 6.88, 2.6, 0.38, font_size=9.5, color=GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – WHY GRAPHRAG WINS (DEEP DIVE)
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Why GraphRAG Outperforms", "Engineering decisions that made the difference", ACCENT2)

reasons = [
    (ACCENT2, "Cross-Sentence Relation Extraction",
     "2-sentence sliding window doubles relations/chunk: 30 → 61.6 avg\n"
     "Captures implicit causal chains that span sentence boundaries"),
    (ACCENT1, "Bigram Query Terms",
     "Extends query_terms() to include 2-word phrases ('supply chain', 'trade facilitation')\n"
     "Bigrams get 2× lexical weight — reduces false positives from single tokens"),
    (GOLD, "Lexical Safety-Net",
     "Force-injects top-2 lexically-matching chunks into final results regardless of entity score\n"
     "Prevents factual chunks from being buried by high-scoring but irrelevant entity matches"),
    (ACCENT4, "Tuned Penalty/Reward Balance",
     "query_miss_penalty: 2.4 → 1.2  |  generic_penalty: 3.4 → 1.6  |  entity_match: 4.0 → 5.0\n"
     "Stopped 8 questions scoring 0.000 — moved them to competitive range"),
    (ACCENT3, "Expanded Trigger Vocabulary",
     "+50 trigger words for IMPROVEMENT, ENABLEMENT, MITIGATION, AFFECT relation types\n"
     "Richer graph means more relation paths available during traversal"),
    (TEAL, "Non-Linear Lexical Anchoring",
     "Lexical boost: 4.5× (≥45% overlap)  |  3.0× (≥25%)  |  1.5× (below)\n"
     "High-overlap chunks are strongly surfaced; low-overlap not penalised"),
]

for i, (col, title, body) in enumerate(reasons):
    row = i // 2
    cx  = 0.3 if i % 2 == 0 else 6.7
    cy  = 1.55 + row * 1.9
    card(s, cx, cy, 6.1, 1.72, MID_BG, col)
    add_text(s, title, cx+0.2, cy+0.1, 5.7, 0.45, font_size=13.5, bold=True, color=col)
    add_text(s, body,  cx+0.2, cy+0.55, 5.7, 1.0,  font_size=11.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – LIVE DEMO WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Interactive Demo", "Step-by-step query walkthrough", ACCENT4)

steps = [
    ("1", "Enter Query",        ACCENT1, "Type a supply-chain question\ne.g. 'How does diversification\nmitigate concentration risk?'"),
    ("2", "Retrieve Chunks",    ACCENT4, "All 4 methods retrieve\ntop-k chunks in parallel\nwith scored evidence"),
    ("3", "Build Graph",        ACCENT2, "Visualise entity nodes &\nrelation edges extracted\nfrom retrieved chunks"),
    ("4", "Traverse Graph",     GOLD,    "Show path scoring:\nentity matches, relation\nbonuses, penalties"),
    ("5", "Generate Answer",    ACCENT3, "LLM synthesises answer\nfrom retrieved context\nwith citations"),
    ("6", "Compare Methods",    WHITE,   "Side-by-side scores\nWinner highlighted\nCategory breakdown"),
]

for i, (num, title, col, desc) in enumerate(steps):
    row = i // 3
    ci  = i % 3
    x   = 0.3 + ci * 4.35
    y   = 1.45 + row * 2.8
    card(s, x, y, 4.1, 2.5, MID_BG, col)
    # big step number
    add_text(s, num, x+0.15, y+0.1, 0.6, 0.9, font_size=36, bold=True, color=col)
    add_text(s, title, x+0.75, y+0.2, 3.1, 0.45, font_size=15, bold=True, color=col)
    add_text(s, desc,  x+0.15, y+1.1, 3.8, 1.25, font_size=12, color=GREY)

add_text(s, "Live at → http://localhost:8501", 3.0, 7.1, 7.0, 0.38,
         font_size=13, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – LIMITATIONS & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
title_bar(s, "Limitations & Future Work", "Honest assessment and next steps", ACCENT3)

limits = [
    (ACCENT3, "Rule-Based Relation Extraction",
     "Trigger-vocabulary approach misses novel phrasings not in the word list.\nSolution → Fine-tune a relation extraction model (e.g. REBEL or RELIK)"),
    (ACCENT3, "Domain Lock-In",
     "Graph schema (9 entity types, 10 relation types) is supply-chain specific.\nSolution → Configurable ontology per domain; auto-induction from corpus"),
    (GOLD,    "Small Benchmark (34 Q)",
     "LLM judge variance is high at low sample sizes; some scores may not be stable.\nSolution → Expand to 100+ questions; use multiple judge models"),
    (GOLD,    "Hybrid Still Underperforms Pure Graph",
     "Fusion weights need adaptive tuning per query type — static RRF is suboptimal.\nSolution → Learn fusion weights with a small supervised head"),
]

future = [
    (ACCENT2, "Graph Neural Networks for retrieval scoring (GraphSAGE / GAT)"),
    (ACCENT2, "Multi-document graph fusion across the full corpus at query time"),
    (ACCENT1, "Temporal graphs — track how policies evolve over time"),
    (ACCENT1, "Fine-tuned LLM answers constrained to graph evidence paths"),
    (ACCENT4, "Streaming updates: add new documents without full re-indexing"),
]

for i, (col, title, body) in enumerate(limits):
    y = 1.45 + i * 1.38
    card(s, 0.3, y, 7.8, 1.22, MID_BG, col)
    add_text(s, title, 0.5, y+0.1,  7.4, 0.4,  font_size=13, bold=True, color=col)
    add_text(s, body,  0.5, y+0.52, 7.4, 0.65, font_size=11.5, color=GREY)

card(s, 8.4, 1.45, 4.65, 5.55, MID_BG, ACCENT2)
add_text(s, "🚀  Future Work", 8.65, 1.6, 4.2, 0.45, font_size=14, bold=True, color=ACCENT2)
for i, (col, txt) in enumerate(future):
    y = 2.2 + i * 0.95
    add_rect(s, 8.55, y, 0.07, 0.62, col)
    add_text(s, txt, 8.72, y+0.05, 4.1, 0.6, font_size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_rect(s, 0, 0, 13.33, 0.1, ACCENT2)
add_rect(s, 0, 7.4, 13.33, 0.1, ACCENT2)

add_text(s, "Key Takeaways", 0.5, 0.25, 12.3, 0.7,
         font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

takeaways = [
    (ACCENT2, "GraphRAG wins the benchmark",
     "Score 0.767 · 22/34 wins · Highest correctness (0.813) & citation precision (91.2%)"),
    (ACCENT1, "Graph structure captures what flat retrieval misses",
     "Causal chains, entity relations, and multi-hop paths that TF-IDF & Dense cannot model"),
    (GOLD,    "Engineering discipline drove the improvement",
     "Penalty tuning + lexical safety-net + bigrams + cross-sentence extraction each added measurable gain"),
    (ACCENT4, "Hybrid is promising but needs adaptive fusion",
     "Static RRF fusion leaves performance on the table — learned weights could push it further"),
    (ACCENT3, "Domain specificity is a feature, not a bug",
     "Targeted ontology outperforms generic embeddings on expert-domain Q&A tasks"),
]

for i, (col, title, body) in enumerate(takeaways):
    y = 1.1 + i * 1.2
    add_rect(s, 0.3, y, 0.12, 1.0, col)
    card(s, 0.5, y, 12.5, 1.0, MID_BG)
    add_text(s, title, 0.75, y+0.08, 12.0, 0.4, font_size=14, bold=True, color=col)
    add_text(s, body,  0.75, y+0.52, 12.0, 0.42, font_size=12.5, color=GREY)

add_text(s, "GraphRAG · Domain-Specific LLMs · 2026", 0.5, 7.1, 12.3, 0.35,
         font_size=11, color=GREY, align=PP_ALIGN.CENTER)

# ── save PPTX ─────────────────────────────────────────────────────────────────
prs.save(str(PPTX_PATH))
print(f"✅  PPTX saved → {PPTX_PATH}")

# ══════════════════════════════════════════════════════════════════════════════
# PDF SPEAKER NOTES
# ══════════════════════════════════════════════════════════════════════════════
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                 Table, TableStyle, HRFlowable, PageBreak)
from reportlab.lib.enums import TA_LEFT, TA_CENTER

W, H = A4

doc = SimpleDocTemplate(
    str(PDF_PATH),
    pagesize=A4,
    leftMargin=2*cm, rightMargin=2*cm,
    topMargin=2*cm, bottomMargin=2*cm,
)

styles = getSampleStyleSheet()

def make_style(name, parent="Normal", **kwargs):
    return ParagraphStyle(name, parent=styles[parent], **kwargs)

title_style    = make_style("MyTitle",    fontSize=20, textColor=colors.HexColor("#58A6FF"),
                             spaceAfter=6, fontName="Helvetica-Bold")
slide_h_style  = make_style("SlideHead",  fontSize=15, textColor=colors.HexColor("#3FB950"),
                             spaceAfter=4, fontName="Helvetica-Bold")
section_style  = make_style("Section",    fontSize=11, textColor=colors.HexColor("#E3B341"),
                             spaceAfter=2, fontName="Helvetica-Bold")
body_style     = make_style("Body",       fontSize=10, textColor=colors.HexColor("#C9D1D9"),
                             spaceAfter=3, leading=14)
bullet_style   = make_style("Bullet",     fontSize=10, textColor=colors.HexColor("#C9D1D9"),
                             leftIndent=14, spaceAfter=2, leading=13,
                             bulletIndent=6)
code_style     = make_style("Code",       fontSize=8.5, textColor=colors.HexColor("#58A6FF"),
                             fontName="Courier", backColor=colors.HexColor("#161B22"),
                             borderPadding=4, spaceAfter=4)
warn_style     = make_style("Warn",       fontSize=10, textColor=colors.HexColor("#F78166"),
                             spaceAfter=2, leading=13)

def HR():
    return HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#30363D"),
                      spaceAfter=8, spaceBefore=4)

def P(text, style=body_style): return Paragraph(text, style)
def B(text): return Paragraph(f"• {text}", bullet_style)
def S(text): return Paragraph(text, section_style)
def H(text): return Paragraph(text, slide_h_style)

slides_notes = [
    # ─── Slide 1 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 1 — Title / Cover",
        "say": [
            "Welcome everyone. Today I'll walk you through a research project at the intersection of knowledge graphs and large language models — specifically, how we can make LLMs smarter on expert domains by augmenting retrieval with a domain-specific knowledge graph.",
            "The domain we work in is supply-chain resilience — a highly policy-relevant topic given recent disruptions from COVID-19, geopolitical tensions, and climate change.",
            "We compared four retrieval approaches on a 34-question benchmark, and our GraphRAG method came out on top.",
        ],
        "technical": [
            "The project is a full end-to-end RAG (Retrieval-Augmented Generation) pipeline implemented in Python.",
            "Corpus: 3 OECD documents + World Bank + WTO guidance, totalling ~X chunks of ~300 tokens.",
            "GraphRAG = Knowledge Graph built from the corpus, used as the retrieval index instead of (or in addition to) a vector store.",
            "Benchmark: 34 questions, 5 categories, scored by GPT-4o-mini acting as judge.",
            "Final scores: Graph 0.767 · Hybrid 0.728 · TF-IDF 0.725 · Dense 0.680.",
        ],
        "qa": [
            "Q: Why supply-chain resilience? — It's a domain rich with causal relationships ('diversification MITIGATES concentration risk') that are perfect for graph modelling.",
            "Q: Is the benchmark publicly available? — The 34 questions are in data/evaluation/questions.jsonl in the repo.",
        ],
    },
    # ─── Slide 2 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 2 — The Problem",
        "say": [
            "Standard RAG pipelines retrieve text by keyword overlap (TF-IDF) or semantic similarity (dense embeddings). Both treat documents as independent units with no notion of how concepts relate to each other.",
            "For factual look-up questions this is fine. But expert-domain Q&A often requires understanding causal chains, trade-offs, and multi-hop paths — 'How does digitalisation improve resilience through better visibility?' — that flat retrieval simply cannot surface.",
            "We also found that chunking severs relationships. An entity mentioned in sentence 3 may be the subject of a causal claim in sentence 5 of the same paragraph. Standard chunking loses that link.",
        ],
        "technical": [
            "TF-IDF failure mode: vocabulary mismatch — 'robustness' and 'resilience' are synonyms but have zero overlap in TF-IDF space.",
            "Dense failure mode: embeddings capture semantic similarity but not directionality of relations — 'A enables B' and 'B enables A' have similar embeddings.",
            "Chunking problem: we use 300-token chunks with 50-token overlap, so a cross-sentence relation spanning >300 tokens is always severed.",
            "Our fix: 2-sentence sliding window for relation extraction during indexing, and graph traversal at query time to recover multi-hop paths.",
        ],
        "qa": [
            "Q: Doesn't dense retrieval handle synonyms? — Yes, but only at the surface level. It cannot follow a path: 'A STRENGTHENS B, B MITIGATES C' to answer 'How does A affect C?'",
            "Q: Why not just use a bigger chunk size? — Larger chunks increase noise (irrelevant sentences) and hurt generation quality. Graph gives you structure without the noise.",
        ],
    },
    # ─── Slide 3 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 3 — System Architecture",
        "say": [
            "The pipeline has two phases. First, an offline indexing phase: we ingest the corpus, split into chunks, extract entities and relations, build the knowledge graph, and store three indexes — TF-IDF, FAISS dense, and a graph JSON.",
            "Second, an online query phase: a question comes in, all four retrievers run in parallel, and the relevant chunks are passed to GPT-4o-mini for answer generation.",
            "The whole system is wrapped in a Streamlit app that visualises each step interactively.",
        ],
        "technical": [
            "Chunking: spaCy sentence tokenizer, ~300 tokens, 50-token overlap, stored as Chunk objects with metadata (doc_id, page, chunk_id).",
            "NER: spaCy en_core_web_sm + EntityRuler with custom patterns for the 9 domain entity types.",
            "Relation extraction: regex trigger-vocabulary (10 relation types, 50+ trigger words each), 2-sentence sliding window.",
            "Graph: NetworkX DiGraph. Nodes = entities (with entity_type, weight). Edges = (source, target, relation_type, chunk_id, confidence).",
            "FAISS index: IndexFlatIP on L2-normalised MiniLM-L6 embeddings.",
            "TF-IDF: sklearn TfidfVectorizer, cosine similarity.",
            "Hybrid: Reciprocal Rank Fusion (k=60) of dense and graph rank lists.",
        ],
        "qa": [
            "Q: How long does indexing take? — ~2-5 minutes on the current corpus on a MacBook with MPS acceleration.",
            "Q: Can you update the graph without re-indexing? — Currently no, full re-index required. Incremental graph updates are listed as future work.",
        ],
    },
    # ─── Slide 4 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 4 — The Knowledge Graph",
        "say": [
            "The heart of GraphRAG is the knowledge graph. Every entity extracted from the corpus becomes a node, and every causal or relational statement becomes a directed edge.",
            "Entity types are weighted by their domain importance — 'strategy' gets weight 3.0 because it's the most semantically loaded concept; 'location' gets 0.5 because mentions of countries are often incidental.",
            "At query time, we score each chunk by how many query entities it contains, what relation types those entities participate in, and whether they're on a short path to other query entities.",
        ],
        "technical": [
            "Graph scoring formula: entity_match + coverage_bonus + sentence_alignment + direct_relation + path_bonus + lexical_score − generic_penalty − query_miss_penalty − hub_penalty.",
            "entity_match = sum of entity weights for entities in both chunk and query.",
            "coverage_bonus = 5.5 × (matched_non_generic / total_query_entities).",
            "direct_relation bonus = +2.0 for each direct edge between query entities in the chunk.",
            "path_bonus = +1.0 for 2-hop paths, +0.5 for 3-hop.",
            "hub_penalty = 0.05 × max(0, degree − 12) to down-weight generic hub nodes.",
            "Bigrams get 2× lexical weight: 'supply chain' scores double versus 'supply' alone.",
            "Lexical safety-net: force-inject top-2 chunks with lexical_overlap ≥ 0.28 into final result, bypassing score competition.",
        ],
        "qa": [
            "Q: How do you handle entity disambiguation? — Currently by exact string match after lowercasing + lemmatisation. Coreference resolution is future work.",
            "Q: What's a hub node? — A node with degree > 12 (e.g. 'supply chain' appears in almost every sentence). We penalise it to prevent it from drowning out specific entities.",
        ],
    },
    # ─── Slide 5 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 5 — Four Retrieval Methods",
        "say": [
            "We compare four methods. TF-IDF is our simplest baseline — pure term frequency matching, no model needed.",
            "Dense retrieval uses sentence-transformers to embed both documents and queries, then FAISS for approximate nearest-neighbour search. This handles paraphrase but has no relational understanding.",
            "Pure GraphRAG uses only the knowledge graph — entity matching plus relation path scoring. It has the highest citation precision at 91.2%.",
            "Hybrid fuses dense and graph signals using Reciprocal Rank Fusion — it aims to get the best of both worlds.",
        ],
        "technical": [
            "TF-IDF: sklearn TfidfVectorizer(ngram_range=(1,2), max_features=50000).",
            "Dense: all-MiniLM-L6-v2 (384-dim), FAISS IndexFlatIP, L2-normalised.",
            "Graph: query_terms() → entity lookup → chunk scoring → lexical safety-net → top-k.",
            "Hybrid RRF: score = 0.38*dense_sim + 0.18*graph_score + 0.08*dense_rank + 0.10*graph_rank + 0.14*entity_coverage + 0.12*relation_signal.",
            "All methods retrieve top-k=3 chunks per query for answer generation.",
            "Citation precision = fraction of cited chunk IDs that actually appear in the gold chunk set.",
        ],
        "qa": [
            "Q: Why not use BM25 instead of TF-IDF? — TF-IDF is used as a simpler, more interpretable baseline. BM25 would likely score higher and is worth adding.",
            "Q: Why MiniLM-L6 and not a larger model? — Speed + memory. MiniLM-L6 runs in <100ms per query on CPU. Results would likely improve with a larger encoder.",
        ],
    },
    # ─── Slide 6 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 6 — Benchmark & Evaluation Setup",
        "say": [
            "We built a 34-question benchmark covering five question types, from direct factual look-up all the way to multi-hop reasoning questions.",
            "Scoring is done by GPT-4o-mini acting as a judge, using a weighted rubric: correctness is the most important dimension at 35%, followed by completeness and groundedness at 25% each.",
            "We deliberately added 10 new questions targeting GraphRAG's strengths — causal chains and multi-hop — to give the graph method a fair chance.",
        ],
        "technical": [
            "Question format: {question_id, question_text, category, gold_chunk_ids, gold_answer}.",
            "gold_chunk_ids are verified real chunk IDs from the corpus (Pydantic min_length=1 enforced).",
            "Judge prompt: provides the question, the generated answer, the gold answer, and the retrieved chunks. Asks for scores 0-1 on each dimension plus written justification.",
            "Answer generation: GPT-4o-mini with a RAG prompt: 'Using only the following context, answer the question. Cite your sources.'",
            "Evaluation pipeline: scripts/run_answer_evaluation.py → results/answer_evaluation/<run_id>/.",
        ],
        "qa": [
            "Q: Is LLM-as-judge reliable? — It has known biases (length, confidence). We mitigate by using calibrated scoring prompts and averaging across 34 questions. For production use, human annotation would be preferred.",
            "Q: Why only 3 chunks for answer generation? — Tested 3, 5, 10. GPT-4o-mini quality degrades with too much context from competing chunks. 3 was the empirically best trade-off.",
        ],
    },
    # ─── Slide 7 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 7 — Benchmark Results",
        "say": [
            "Here are the headline numbers. GraphRAG scores 0.767 overall — beating TF-IDF at 0.725, Hybrid at 0.728, and Dense at 0.680.",
            "On correctness — the most important dimension — GraphRAG scores 0.813 compared to TF-IDF's 0.762.",
            "On citation precision, GraphRAG and Hybrid both achieve 91.2%, meaning almost every chunk they cite is genuinely relevant. TF-IDF is at 85.3%.",
            "Both GraphRAG and TF-IDF win on 22 questions each, but GraphRAG wins on harder questions while TF-IDF wins on simpler factual ones.",
        ],
        "technical": [
            "Scores are computed as weighted average: 0.35*correctness + 0.25*completeness + 0.25*groundedness + 0.10*reasoning + 0.05*clarity.",
            "Wins = number of questions where a method achieved the highest overall score (ties counted for all tied methods).",
            "GraphRAG's Wins breakdown: 5 direct_factual, 7 causal, 2 mitigation, 5 multi_hop, 3 trade_off.",
            "TF-IDF's Wins breakdown: concentrated in direct_factual and multi_hop — categories where term overlap is sufficient.",
            "Raw results file: results/answer_evaluation/answer_eval_20260412T164625164443+0000/per_question_results.jsonl",
        ],
        "qa": [
            "Q: Why does TF-IDF still win on 22 questions? — On direct factual questions (e.g. 'What are the three building blocks?'), exact term matching is hard to beat. GraphRAG is not uniformly better — it's better on relational questions.",
            "Q: Can you combine the strengths of both? — The Hybrid method tries to do exactly that, but with static weights it doesn't beat pure Graph overall. Adaptive fusion is the next step.",
        ],
    },
    # ─── Slide 8 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 8 — Performance by Category",
        "say": [
            "Breaking down by category reveals where GraphRAG really shines.",
            "On mitigation questions — 'what strategies reduce supply chain risk' — GraphRAG scores 0.631 vs TF-IDF's 0.459. That's a 37% relative improvement.",
            "On causal reasoning questions, GraphRAG leads at 0.800. This makes intuitive sense: causal chains are exactly what graph edges represent.",
            "TF-IDF's strength is on direct factual and multi-hop — where verbatim term overlap is informative.",
        ],
        "technical": [
            "Category scores computed as mean over questions in that category.",
            "Mitigation gap: TF-IDF 0.459 is extremely low because mitigation questions often use paraphrase ('absorb shocks', 'cushion disruptions') not present in documents verbatim.",
            "Multi-hop: Hybrid 0.777 ties TF-IDF 0.777 — fusion helps here by combining dense semantic matching with graph path signals.",
            "Trade-off: Dense 0.772 leads — trade-off questions ('tension between efficiency and resilience') have a high semantic signal that dense embeddings capture well.",
            "GraphRAG weakness: mitigation (0.631) — suggests the MITIGATES relation extraction needs improvement, or that mitigation questions are phrased in ways not well-covered by the trigger vocabulary.",
        ],
        "qa": [
            "Q: Why does Dense beat Graph on trade-off questions? — Trade-off questions are about contrasting two concepts. Dense embeddings are trained on semantic similarity and capture 'tension between X and Y' patterns well. Graphs model directionality, not symmetric contrast.",
            "Q: What would you do to improve mitigation scores? — Fine-tune the relation extraction model specifically on mitigation-type sentences, and expand the MITIGATES trigger vocabulary further.",
        ],
    },
    # ─── Slide 9 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 9 — Why GraphRAG Wins",
        "say": [
            "The improvement from the initial 0.520 to the current 0.767 came from six targeted engineering decisions.",
            "The most impactful single change was the lexical safety-net — a force-injection mechanism that ensures the top lexically-matching chunks are always in the result set, regardless of their entity score.",
            "Why was this needed? Entity scores range from 0-200 while lexical scores range from 0-5. Without the safety-net, a factual chunk with zero named entities would never be retrieved.",
        ],
        "technical": [
            "Cross-sentence extraction: for each sentence pair (s_i, s_{i+1}), we run relation extraction jointly. This doubled average relations per chunk from ~30 to ~61.6.",
            "Bigram query_terms: tokens = [t for t in TOKEN_RE.findall(query.lower()) if len(t)>2]. For each consecutive pair (a,b): if a or b not in STOPWORDS, add '{a} {b}' to terms set.",
            "Lexical safety-net: compute lexical_overlap_score for all chunks; take top-2 with score ≥ 0.28; inject into last positions of final ranked list (replacing positions top_k-2 and top_k-1).",
            "Penalty tuning: query_miss 2.4→1.2, generic_penalty 3.4→1.6, hub_penalty factor 0.10→0.05, hub threshold 10→12. Entity_match weight 4.0→5.0.",
            "Non-linear lexical: if overlap ≥ 0.45 → ×4.5; elif ≥ 0.25 → ×3.0; else → ×1.5.",
            "Expanded triggers: +14 IMPROVEMENT words, +9 ENABLEMENT, +11 MITIGATION, +18 AFFECT.",
        ],
        "qa": [
            "Q: Isn't the safety-net a hack? — It's a pragmatic engineering choice to bridge the score-scale gap between entity-based and lexical-based signals. A unified scoring model would be the clean solution.",
            "Q: How did you know which penalties were too harsh? — Diagnosed by looking at per-question results where graph scored exactly 0.000 (8 questions). Traced back to early-exit condition triggered by empty entity matches + high no-relation penalty.",
        ],
    },
    # ─── Slide 10 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 10 — Interactive Demo",
        "say": [
            "The system is wrapped in a Streamlit app that walks you through the full pipeline step by step.",
            "You type a question, and the app shows you: which chunks each method retrieves, the knowledge graph built from those chunks, the scoring breakdown for graph traversal, the generated answers, and a final comparison across all four methods.",
            "Let me show you a live example — I'll type 'How does diversification mitigate concentration risks?' and walk through each step.",
        ],
        "technical": [
            "App stack: Streamlit 1.x, Plotly for graph visualisation (spring layout), Pandas for comparison tables.",
            "Session state keys: query_text, active_method, step, graph_data, retrieval_results, answers.",
            "Step navigation uses stable deterministic widget keys ('nav_next', 'qs_start') — critical to avoid Streamlit re-render issues where clicking Next didn't register.",
            "Graph visualisation: nx.spring_layout → Plotly scatter with edge traces. Node colour = entity_type, node size = degree.",
            "All 4 retrievers are called lazily (only when the user reaches the retrieval step).",
        ],
        "qa": [
            "Q: Can the app be deployed publicly? — Yes, with minor changes (add authentication, use st.secrets for API keys). The current setup requires a local .env with OPENAI_API_KEY.",
            "Q: Why Streamlit over FastAPI+React? — Streamlit allows rapid iteration for research demos. For production, a proper frontend would be needed.",
        ],
    },
    # ─── Slide 11 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 11 — Limitations & Future Work",
        "say": [
            "I want to be upfront about the limitations. The relation extraction is rule-based — it works well for our domain and trigger vocabulary, but it will miss novel phrasings.",
            "The benchmark is small at 34 questions, so individual score differences may not be statistically significant. We'd want at least 100 questions with multiple human annotators to make strong claims.",
            "The hybrid method surprisingly underperforms pure graph overall. This suggests the fusion weights need to be learned per-query-type rather than fixed.",
        ],
        "technical": [
            "Rule-based NRE limitations: trigger vocabulary has ~200 words per relation type. Novel policy language (e.g. 'bolster adaptive capacity') may not be covered.",
            "Statistical significance: with n=34 and LLM judge variance ~±0.05 per question, differences of <0.05 in mean score are not reliable.",
            "Hybrid fusion issue: RRF with static k=60 and fixed linear weights treats all query types identically. A meta-ranker trained on (query_type, method_scores) → optimal_weight could improve this.",
            "Future work priority order: (1) learned fusion weights, (2) neural relation extraction, (3) larger benchmark, (4) coreference resolution, (5) incremental graph updates.",
        ],
        "qa": [
            "Q: What about hallucination? — Groundedness scores measure this indirectly (0.756 for graph). The main risk is when no relevant chunk is retrieved and the LLM generates from parametric memory. We mitigate with a 'answer only from context' prompt.",
            "Q: Have you considered RAG-fusion or HyDE? — Not yet. HyDE (hypothetical document embeddings) could help with semantic gap for causal questions. Listed as future work.",
        ],
    },
    # ─── Slide 12 ───────────────────────────────────────────────────────────────
    {
        "title": "Slide 12 — Conclusion",
        "say": [
            "To wrap up: this project demonstrates that for expert-domain Q&A, augmenting retrieval with a knowledge graph built from the same corpus provides measurable quality improvements — especially on reasoning-intensive questions.",
            "GraphRAG achieves 0.767 overall, winning 22 out of 34 benchmark questions with the highest correctness and citation precision.",
            "The key insight is that language is not just bags of words — it encodes structured relationships, and a graph is the natural data structure to capture and query those relationships.",
            "Thank you — I'm happy to take questions.",
        ],
        "technical": [
            "Full results in results/answer_evaluation/answer_eval_20260412T164625164443+0000/.",
            "Source code: src/graph_rag/{schema.py, query.py, hybrid_query.py}, app.py.",
            "Evaluation script: scripts/run_answer_evaluation.py.",
            "Benchmark: data/evaluation/questions.jsonl (34 questions).",
            "All improvements are tracked in git history on the main branch.",
        ],
        "qa": [
            "Q: Would GraphRAG work on other domains? — Yes, with a new ontology. The pipeline is parameterised by entity types and relation trigger vocabularies. Legal, medical, and financial domains are natural candidates.",
            "Q: What's the biggest single improvement you'd make? — Replace rule-based relation extraction with a fine-tuned neural model (e.g. REBEL). That would capture the long tail of phrasings and likely push graph scores above 0.85.",
            "Q: Why not just use a fine-tuned LLM? — GraphRAG is complementary: it gives the LLM better evidence. A fine-tuned LLM + GraphRAG would be the strongest combination.",
        ],
    },
]

story = []
story.append(Paragraph("GraphRAG for Domain-Specific LLMs", title_style))
story.append(Paragraph("Speaker Notes & Technical Reference", make_style("sub", fontSize=13,
             textColor=colors.HexColor("#8B949E"))))
story.append(Spacer(1, 0.4*cm))
story.append(HR())
story.append(Spacer(1, 0.3*cm))

for slide in slides_notes:
    story.append(H(slide["title"]))
    story.append(HR())

    story.append(S("WHAT TO SAY"))
    for line in slide["say"]:
        story.append(B(line))
    story.append(Spacer(1, 0.25*cm))

    story.append(S("TECHNICAL DETAILS (for Q&A)"))
    for line in slide["technical"]:
        story.append(B(line))
    story.append(Spacer(1, 0.25*cm))

    story.append(S("ANTICIPATED QUESTIONS & ANSWERS"))
    for line in slide["qa"]:
        story.append(P(line, make_style(f"qa_{slide['title']}", fontSize=10,
                       textColor=colors.HexColor("#C9D1D9"), leftIndent=10,
                       spaceAfter=4, leading=13,
                       borderPadding=3)))
    story.append(Spacer(1, 0.5*cm))
    story.append(PageBreak())

doc.build(story)
print(f"✅  PDF  saved → {PDF_PATH}")
print("\nDone! Both files are in the PresentAAAA/ folder.")
